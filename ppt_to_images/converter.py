"""
Core conversion logic for PPT/PPTX/PDF to images.
"""

import os
import subprocess
import logging
from pathlib import Path
from typing import Union, Optional, List, Dict, Any, Literal, BinaryIO

from pdf2image import convert_from_path
from PIL import Image

from .utils import (
    FileType, 
    detect_file_type, 
    image_to_base64, 
    TempFileManager,
    ensure_directory
)

logger = logging.getLogger(__name__)


class ConversionError(Exception):
    """Exception raised when conversion fails."""
    pass


class PPTConverter:
    """
    PPT/PPTX/PDF to image sequence converter.
    
    Example:
        >>> converter = PPTConverter(dpi=200)
        >>> result = converter.convert("input.pptx", output_dir="./images")
        >>> print(result["count"])
        10
    """
    
    def __init__(
        self, 
        dpi: int = 200, 
        cleanup_temp: bool = True,
        temp_dir: Optional[str] = None
    ):
        """
        Initialize the converter.
        
        Args:
            dpi: Image DPI (resolution), default 200
            cleanup_temp: Whether to clean up temporary files
            temp_dir: Custom temporary directory
        """
        self.dpi = dpi
        self.cleanup_temp = cleanup_temp
        self.temp_dir = temp_dir
    
    def convert(
        self,
        input_path: Union[str, Path, BinaryIO],
        output_dir: Optional[Union[str, Path]] = None,
        output_format: Literal["file", "base64", "both"] = "file",
        extract_notes: bool = False
    ) -> Dict[str, Any]:
        """
        Convert PPT/PPTX/PDF to image sequence.
        
        Args:
            input_path: Input file path or file-like object
            output_dir: Output directory (required for "file" and "both" formats)
            output_format: Output format - "file", "base64", or "both"
            extract_notes: Whether to extract notes (备注) from slides
            
        Returns:
            Dictionary containing:
                - images: List of file paths or base64 strings
                - images_base64: List of base64 strings (only for "both" format)
                - count: Number of images
                - texts: List of notes content (if extract_notes=True)
                - format: Output format used
                
        Raises:
            ConversionError: If conversion fails
            ValueError: If invalid arguments provided
        """
        # Validate arguments
        if output_format in ("file", "both") and not output_dir:
            raise ValueError(f"output_dir is required for format '{output_format}'")
        
        # Handle file path
        if isinstance(input_path, (str, Path)):
            input_path = Path(input_path)
            if not input_path.exists():
                raise FileNotFoundError(f"Input file not found: {input_path}")
        
        # Create output directory
        if output_dir:
            output_dir = ensure_directory(output_dir)
        
        with TempFileManager(cleanup=self.cleanup_temp, temp_dir=self.temp_dir) as temp:
            # Handle file-like object
            if hasattr(input_path, 'read'):
                content = input_path.read()
                filename = getattr(input_path, 'name', 'input.pptx')
                input_path = Path(temp.save_uploaded_file(content, Path(filename).name))
            
            # Detect file type
            file_type = detect_file_type(input_path)
            if file_type == FileType.UNKNOWN:
                raise ConversionError(f"Unsupported file type: {input_path.suffix}")
            
            # Convert to PDF if needed
            pdf_path = self._ensure_pdf(input_path, file_type, temp)
            
            # Convert PDF to images
            images = self._pdf_to_images(pdf_path)
            
            # Extract notes if requested
            texts = []
            if extract_notes and file_type in (FileType.PPT, FileType.PPTX):
                pptx_path = self._ensure_pptx(input_path, file_type, temp)
                texts = self._extract_text(pptx_path)
            
            # Process output
            result = self._process_output(
                images, output_dir, output_format, texts, extract_notes
            )
            
            return result
    
    def _ensure_pdf(
        self, 
        input_path: Path, 
        file_type: FileType, 
        temp: TempFileManager
    ) -> Path:
        """
        Ensure we have a PDF file for conversion.
        
        Args:
            input_path: Original input file path
            file_type: Detected file type
            temp: Temporary file manager
            
        Returns:
            Path to PDF file
        """
        if file_type == FileType.PDF:
            return input_path
        
        # Convert PPT to PPTX first if needed
        if file_type == FileType.PPT:
            pptx_path = self._convert_ppt_to_pptx(input_path, temp)
        else:
            pptx_path = input_path
        
        # Convert PPTX to PDF
        pdf_path = self._convert_pptx_to_pdf(pptx_path, temp)
        return pdf_path
    
    def _ensure_pptx(
        self,
        input_path: Path,
        file_type: FileType,
        temp: TempFileManager
    ) -> Path:
        """
        Ensure we have a PPTX file for text extraction.
        
        Args:
            input_path: Original input file path
            file_type: Detected file type
            temp: Temporary file manager
            
        Returns:
            Path to PPTX file
        """
        if file_type == FileType.PPTX:
            return input_path
        elif file_type == FileType.PPT:
            return self._convert_ppt_to_pptx(input_path, temp)
        else:
            raise ConversionError("Cannot extract text from PDF files")
    
    def _convert_ppt_to_pptx(self, ppt_path: Path, temp: TempFileManager) -> Path:
        """
        Convert PPT to PPTX using LibreOffice.
        
        Args:
            ppt_path: Path to PPT file
            temp: Temporary file manager
            
        Returns:
            Path to converted PPTX file
        """
        logger.info(f"Converting PPT to PPTX: {ppt_path}")
        
        output_dir = temp._created_dir
        
        try:
            result = subprocess.run(
                [
                    "soffice", "--headless", "--convert-to", "pptx",
                    "--outdir", output_dir, str(ppt_path)
                ],
                capture_output=True,
                text=True,
                timeout=300
            )
            
            if result.returncode != 0:
                raise ConversionError(f"LibreOffice conversion failed: {result.stderr}")
            
            pptx_path = Path(output_dir) / f"{ppt_path.stem}.pptx"
            if not pptx_path.exists():
                raise ConversionError("PPTX file not created after conversion")
            
            return pptx_path
            
        except subprocess.TimeoutExpired:
            raise ConversionError("Conversion timed out")
        except FileNotFoundError:
            raise ConversionError("LibreOffice not found. Please install LibreOffice.")
    
    def _convert_pptx_to_pdf(self, pptx_path: Path, temp: TempFileManager) -> Path:
        """
        Convert PPTX to PDF using LibreOffice.
        
        Args:
            pptx_path: Path to PPTX file
            temp: Temporary file manager
            
        Returns:
            Path to converted PDF file
        """
        logger.info(f"Converting PPTX to PDF: {pptx_path}")
        
        output_dir = temp._created_dir
        
        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", output_dir, str(pptx_path)
                ],
                capture_output=True,
                text=True,
                timeout=300
            )
            
            if result.returncode != 0:
                raise ConversionError(f"LibreOffice conversion failed: {result.stderr}")
            
            pdf_path = Path(output_dir) / f"{pptx_path.stem}.pdf"
            if not pdf_path.exists():
                raise ConversionError("PDF file not created after conversion")
            
            return pdf_path
            
        except subprocess.TimeoutExpired:
            raise ConversionError("Conversion timed out")
        except FileNotFoundError:
            raise ConversionError("LibreOffice not found. Please install LibreOffice.")
    
    def _pdf_to_images(self, pdf_path: Path) -> List[Image.Image]:
        """
        Convert PDF to PIL Image objects.
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of PIL Image objects
        """
        logger.info(f"Converting PDF to images: {pdf_path}")
        
        try:
            images = convert_from_path(pdf_path, dpi=self.dpi)
            return images
        except Exception as e:
            raise ConversionError(f"PDF to image conversion failed: {e}")
    
    def _extract_text(self, pptx_path: Path) -> List[str]:
        """
        Extract notes (备注) content from PPTX file.
        Only extracts notes, not slide text content (which will be extracted via OCR later).
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            List of notes content per slide
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed, skipping text extraction")
            return []
        
        logger.info(f"Extracting notes from: {pptx_path}")
        
        texts = []
        try:
            presentation = Presentation(str(pptx_path))
            
            for slide in presentation.slides:
                # Only extract notes, not slide text content
                note = ""
                if slide.has_notes_slide:
                    note = slide.notes_slide.notes_text_frame.text
                
                texts.append(note)
            
            return texts
            
        except Exception as e:
            logger.warning(f"Notes extraction failed: {e}")
            return []
    
    def _process_output(
        self,
        images: List[Image.Image],
        output_dir: Optional[Path],
        output_format: str,
        texts: List[str],
        extract_notes: bool
    ) -> Dict[str, Any]:
        """
        Process and format output.
        
        Args:
            images: List of PIL Image objects
            output_dir: Output directory
            output_format: Output format
            texts: Extracted notes content
            extract_notes: Whether notes extraction was requested
            
        Returns:
            Result dictionary
        """
        result = {
            "count": len(images),
            "format": output_format
        }
        
        if output_format == "file":
            result["images"] = self._save_images(images, output_dir)
        elif output_format == "base64":
            result["images"] = self._images_to_base64(images)
        elif output_format == "both":
            result["images"] = self._save_images(images, output_dir)
            result["images_base64"] = self._images_to_base64(images)
        
        if extract_notes:
            result["texts"] = texts
        
        return result
    
    def _save_images(self, images: List[Image.Image], output_dir: Path) -> List[str]:
        """
        Save images to files.
        
        Args:
            images: List of PIL Image objects
            output_dir: Output directory
            
        Returns:
            List of file paths
        """
        paths = []
        for i, image in enumerate(images, start=1):
            path = output_dir / f"{i}.png"
            image.save(str(path), "PNG")
            paths.append(str(path))
            logger.debug(f"Saved image: {path}")
        
        return paths
    
    def _images_to_base64(self, images: List[Image.Image]) -> List[str]:
        """
        Convert images to base64 strings.
        
        Args:
            images: List of PIL Image objects
            
        Returns:
            List of base64 encoded strings
        """
        import io
        
        base64_list = []
        for image in images:
            buffer = io.BytesIO()
            image.save(buffer, format="PNG")
            base64_str = image_to_base64_from_bytes(buffer.getvalue())
            base64_list.append(base64_str)
        
        return base64_list


def image_to_base64_from_bytes(data: bytes) -> str:
    """
    Convert image bytes to base64 string.
    
    Args:
        data: Image data as bytes
        
    Returns:
        Base64 encoded string
    """
    import base64
    return base64.b64encode(data).decode('utf-8')

